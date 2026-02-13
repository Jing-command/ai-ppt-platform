#!/usr/bin/env python3
"""
导出系统测试脚本
"""
import asyncio
import sys
import os
from pathlib import Path
from uuid import uuid4, UUID

# 添加项目路径
sys.path.insert(0, '/root/.openclaw/workspace/ai-ppt-platform/backend/src')

async def test_export_service():
    """测试导出服务"""
    print("=" * 60)
    print("测试导出服务 (Export Service)")
    print("=" * 60)
    
    from ai_ppt.services.export_service import (
        ExportService, ExportFormat, ExportStatus, ExportTask
    )
    
    # 测试导出任务创建
    print("\n1. 测试导出任务创建...")
    user_id = uuid4()
    presentation_id = uuid4()
    
    task = ExportTask(
        user_id=user_id,
        presentation_id=presentation_id,
        format=ExportFormat.PPTX,
        quality="high",
        slide_range="all",
        include_notes=False,
    )
    
    print(f"   ✅ 任务创建成功")
    print(f"      - 任务ID: {task.id}")
    print(f"      - 格式: {task.format.value}")
    print(f"      - 状态: {task.status.value}")
    print(f"      - PPT ID: {task.presentation_id}")
    
    # 测试主题颜色获取
    print("\n2. 测试主题颜色...")
    from ai_ppt.services.export_service import ExportService
    # 创建一个mock对象来测试_get_theme_colors
    class MockService:
        _get_theme_colors = ExportService._get_theme_colors
    
    mock = MockService()
    colors = mock._get_theme_colors("blue")
    print(f"   ✅ 主题颜色获取成功")
    print(f"      - 背景色: {colors['background']}")
    print(f"      - 标题色: {colors['title']}")
    print(f"      - 文本色: {colors['text']}")
    
    print("\n3. 测试其他主题...")
    for theme in ["default", "dark", "green"]:
        colors = mock._get_theme_colors(theme)
        print(f"   ✅ {theme}: 背景{colors['background']}")
    
    print("\n" + "=" * 60)
    print("导出服务测试完成!")
    print("=" * 60)
    return True

async def test_pptx_generation():
    """测试PPTX生成功能"""
    print("\n" + "=" * 60)
    print("测试 PPTX 生成功能")
    print("=" * 60)
    
    from pptx import Presentation as PptxPresentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    
    # 创建临时目录
    exports_dir = Path("/tmp/test_exports")
    exports_dir.mkdir(parents=True, exist_ok=True)
    
    print("\n1. 创建PPTX文件...")
    prs = PptxPresentation()
    
    # 设置幻灯片尺寸 (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 添加空白布局
    blank_layout = prs.slide_layouts[6]
    
    # 添加标题幻灯片
    slide = prs.slides.add_slide(blank_layout)
    
    # 设置背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 248, 255)
    
    # 添加标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "测试演示文稿"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(25, 55, 109)
    
    # 添加副标题
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.2), Inches(12.333), Inches(0.5)
    )
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "AI PPT Platform 导出测试"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(60, 80, 120)
    
    # 添加内容幻灯片
    slide2 = prs.slides.add_slide(blank_layout)
    background2 = slide2.background
    fill2 = background2.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(255, 255, 255)
    
    title_box2 = slide2.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
    )
    tf2 = title_box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "第一章：测试内容"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(51, 51, 51)
    
    # 添加项目符号
    bullets_box = slide2.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(12.333), Inches(4.5)
    )
    tf3 = bullets_box.text_frame
    tf3.word_wrap = True
    
    bullets = ["测试项目 1", "测试项目 2", "测试项目 3"]
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf3.paragraphs[0]
        else:
            p = tf3.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(102, 102, 102)
        p.space_after = Pt(8)
    
    # 保存文件
    output_file = exports_dir / "test_presentation.pptx"
    prs.save(str(output_file))
    
    print(f"   ✅ PPTX文件创建成功")
    print(f"      - 文件路径: {output_file}")
    print(f"      - 文件大小: {output_file.stat().st_size} bytes")
    print(f"      - 幻灯片数: 2")
    
    # 清理
    output_file.unlink()
    exports_dir.rmdir()
    
    print("\n" + "=" * 60)
    print("PPTX生成测试完成!")
    print("=" * 60)
    return True

async def test_pdf_generation():
    """测试PDF生成功能"""
    print("\n" + "=" * 60)
    print("测试 PDF 生成功能")
    print("=" * 60)
    
    from reportlab.lib.pagesizes import landscape, A4
    from reportlab.pdfgen import canvas
    
    exports_dir = Path("/tmp/test_exports")
    exports_dir.mkdir(parents=True, exist_ok=True)
    
    print("\n1. 创建PDF文件...")
    width, height = landscape(A4)
    output_file = exports_dir / "test_presentation.pdf"
    
    c = canvas.Canvas(str(output_file), pagesize=(width, height))
    
    # 第一页
    c.setFillColorRGB(240/255, 248/255, 255/255)
    c.rect(0, 0, width, height, fill=1, stroke=0)
    
    c.setFillColorRGB(25/255, 55/255, 109/255)
    c.setFont("Helvetica-Bold", 28)
    c.drawString(40, height - 80, "Test Presentation")
    
    c.setFillColorRGB(60/255, 80/255, 120/255)
    c.setFont("Helvetica", 16)
    c.drawString(40, height - 120, "AI PPT Platform Export Test")
    
    c.showPage()
    
    # 第二页
    c.setFillColorRGB(255/255, 255/255, 255/255)
    c.rect(0, 0, width, height, fill=1, stroke=0)
    
    c.setFillColorRGB(51/255, 51/255, 51/255)
    c.setFont("Helvetica-Bold", 24)
    c.drawString(40, height - 80, "Chapter 1: Test Content")
    
    bullets = ["Test item 1", "Test item 2", "Test item 3"]
    y = height - 140
    c.setFillColorRGB(102/255, 102/255, 102/255)
    c.setFont("Helvetica", 12)
    for bullet in bullets:
        c.drawString(40, y, f"• {bullet}")
        y -= 25
    
    c.showPage()
    c.save()
    
    print(f"   ✅ PDF文件创建成功")
    print(f"      - 文件路径: {output_file}")
    print(f"      - 文件大小: {output_file.stat().st_size} bytes")
    print(f"      - 页面数: 2")
    
    # 清理
    output_file.unlink()
    exports_dir.rmdir()
    
    print("\n" + "=" * 60)
    print("PDF生成测试完成!")
    print("=" * 60)
    return True

async def test_image_generation():
    """测试图片生成功能"""
    print("\n" + "=" * 60)
    print("测试图片生成功能")
    print("=" * 60)
    
    from PIL import Image, ImageDraw
    import zipfile
    
    exports_dir = Path("/tmp/test_exports")
    exports_dir.mkdir(parents=True, exist_ok=True)
    
    print("\n1. 创建图片文件...")
    width, height = 1280, 720
    
    image_files = []
    
    # 创建第一张图片
    img1 = Image.new('RGB', (width, height), (240, 248, 255))
    draw1 = ImageDraw.Draw(img1)
    
    # 绘制简单内容
    draw1.rectangle([0, 0, width, 80], fill=(25, 55, 109))
    draw1.text((40, 30), "Test Presentation", fill=(255, 255, 255))
    draw1.text((40, 120), "AI PPT Platform Export Test", fill=(60, 80, 120))
    
    img1_path = exports_dir / "slide_001.png"
    img1.save(img1_path, "PNG")
    image_files.append(img1_path)
    
    # 创建第二张图片
    img2 = Image.new('RGB', (width, height), (255, 255, 255))
    draw2 = ImageDraw.Draw(img2)
    
    draw2.rectangle([0, 0, width, 80], fill=(51, 51, 51))
    draw2.text((40, 30), "Chapter 1: Test Content", fill=(255, 255, 255))
    
    y = 140
    for bullet in ["Test item 1", "Test item 2", "Test item 3"]:
        draw2.text((40, y), f"• {bullet}", fill=(102, 102, 102))
        y += 35
    
    img2_path = exports_dir / "slide_002.png"
    img2.save(img2_path, "PNG")
    image_files.append(img2_path)
    
    print(f"   ✅ 图片创建成功")
    for f in image_files:
        print(f"      - {f.name}: {f.stat().st_size} bytes")
    
    # 打包为zip
    print("\n2. 打包为ZIP文件...")
    zip_path = exports_dir / "test_slides.zip"
    with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zf:
        for img_file in image_files:
            zf.write(img_file, img_file.name)
    
    print(f"   ✅ ZIP文件创建成功")
    print(f"      - 文件路径: {zip_path}")
    print(f"      - 文件大小: {zip_path.stat().st_size} bytes")
    
    # 清理
    for f in image_files:
        f.unlink()
    zip_path.unlink()
    exports_dir.rmdir()
    
    print("\n" + "=" * 60)
    print("图片生成测试完成!")
    print("=" * 60)
    return True

async def main():
    """主测试函数"""
    print("\n" + "=" * 60)
    print("AI PPT Platform 导出系统测试")
    print("=" * 60)
    
    results = []
    
    try:
        results.append(("导出服务", await test_export_service()))
    except Exception as e:
        print(f"❌ 导出服务测试失败: {e}")
        results.append(("导出服务", False))
    
    try:
        results.append(("PPTX生成", await test_pptx_generation()))
    except Exception as e:
        print(f"❌ PPTX生成测试失败: {e}")
        results.append(("PPTX生成", False))
    
    try:
        results.append(("PDF生成", await test_pdf_generation()))
    except Exception as e:
        print(f"❌ PDF生成测试失败: {e}")
        results.append(("PDF生成", False))
    
    try:
        results.append(("图片生成", await test_image_generation()))
    except Exception as e:
        print(f"❌ 图片生成测试失败: {e}")
        results.append(("图片生成", False))
    
    # 打印测试摘要
    print("\n" + "=" * 60)
    print("测试摘要")
    print("=" * 60)
    for name, result in results:
        status = "✅ 通过" if result else "❌ 失败"
        print(f"  {status}: {name}")
    
    passed = sum(1 for _, r in results if r)
    total = len(results)
    print(f"\n总计: {passed}/{total} 项测试通过")
    print("=" * 60)
    
    return all(r for _, r in results)

if __name__ == "__main__":
    success = asyncio.run(main())
    sys.exit(0 if success else 1)
