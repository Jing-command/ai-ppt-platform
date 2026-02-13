"""
导出服务 - 处理PPT导出为各种格式
支持 PPTX、PDF、PNG/JPG 导出
"""

from __future__ import annotations

import os
import uuid
from datetime import datetime, timedelta
from enum import Enum as PyEnum
from pathlib import Path
from typing import TYPE_CHECKING, Optional, Union
from uuid import UUID

from PIL import Image, ImageDraw, ImageFont
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

if TYPE_CHECKING:
    from ai_ppt.domain.models.presentation import Presentation
    from ai_ppt.domain.models.slide import Slide


class ExportFormat(str, PyEnum):
    """导出格式枚举"""

    PPTX = "pptx"
    PDF = "pdf"
    PNG = "png"
    JPG = "jpg"


class ExportStatus(str, PyEnum):
    """导出状态枚举"""

    PENDING = "pending"
    PROCESSING = "processing"
    COMPLETED = "completed"
    FAILED = "failed"


class ExportTask:
    """导出任务模型（内存存储）"""

    def __init__(
        self,
        user_id: UUID,
        presentation_id: UUID,
        format: ExportFormat,
        quality: str = "standard",
        slide_range: Optional[str] = None,
        include_notes: bool = False,
    ):
        self.id = uuid.uuid4()
        self.user_id = user_id
        self.presentation_id = presentation_id
        self.format = format
        self.quality = quality
        self.slide_range = slide_range
        self.include_notes = include_notes
        self.status = ExportStatus.PENDING
        self.progress = 0
        self.file_path: Optional[str] = None
        self.file_size: Optional[int] = None
        self.error_message: Optional[str] = None
        self.created_at = datetime.utcnow()
        self.completed_at: Optional[datetime] = None
        self.expires_at: Optional[datetime] = None


# 内存任务存储（生产环境应使用Redis + Celery）
_export_tasks: dict[UUID, ExportTask] = {}


class ExportServiceError(Exception):
    """导出服务错误基类"""


class ExportNotFoundError(ExportServiceError):
    """导出任务不存在"""


class ExportFailedError(ExportServiceError):
    """导出失败错误"""


class ExportService:
    """
    导出服务

    处理 PPT 导出为多种格式：
    - PPTX: PowerPoint 格式
    - PDF: PDF 文档
    - PNG/JPG: 图片格式
    """

    def __init__(self, session: AsyncSession, exports_dir: str = "exports"):
        """
        初始化导出服务

        Args:
            session: 数据库会话
            exports_dir: 导出文件存储目录
        """
        self._session = session
        self._exports_dir = Path(exports_dir)
        self._exports_dir.mkdir(parents=True, exist_ok=True)

    async def create_task(
        self,
        user_id: UUID,
        presentation_id: UUID,
        format: ExportFormat,
        quality: str = "standard",
        slide_range: Optional[str] = None,
        include_notes: bool = False,
    ) -> ExportTask:
        """
        创建导出任务

        Args:
            user_id: 用户ID
            presentation_id: PPT ID
            format: 导出格式
            quality: 导出质量
            slide_range: 页面范围
            include_notes: 是否包含备注

        Returns:
            导出任务
        """
        task = ExportTask(
            user_id=user_id,
            presentation_id=presentation_id,
            format=format,
            quality=quality,
            slide_range=slide_range,
            include_notes=include_notes,
        )
        _export_tasks[task.id] = task
        return task

    async def get_task(self, task_id: UUID, user_id: UUID) -> Optional[ExportTask]:
        """
        获取导出任务

        Args:
            task_id: 任务ID
            user_id: 用户ID（用于权限检查）

        Returns:
            导出任务，不存在或无权限返回 None
        """
        task = _export_tasks.get(task_id)
        if task and task.user_id == user_id:
            return task
        return None

    async def process_export(self, task_id: UUID) -> None:
        """
        处理导出任务（异步执行）

        Args:
            task_id: 任务ID
        """
        task = _export_tasks.get(task_id)
        if not task:
            return

        try:
            task.status = ExportStatus.PROCESSING
            task.progress = 10

            # 获取 PPT 和幻灯片
            presentation = await self._get_presentation(task.presentation_id)
            if not presentation:
                raise ExportFailedError("Presentation not found")

            task.progress = 30
            slides = await self._get_slides(task.presentation_id)

            # 根据范围过滤幻灯片
            if task.slide_range and task.slide_range != "all":
                slides = self._filter_slides_by_range(slides, task.slide_range)

            task.progress = 50

            # 根据格式导出
            if task.format == ExportFormat.PPTX:
                file_path = await self._export_pptx(presentation, slides, task)
            elif task.format == ExportFormat.PDF:
                file_path = await self._export_pdf(presentation, slides, task)
            elif task.format in (ExportFormat.PNG, ExportFormat.JPG):
                file_path = await self._export_images(presentation, slides, task)
            else:
                raise ExportFailedError(f"Unsupported format: {task.format}")

            task.file_path = file_path
            task.file_size = os.path.getsize(file_path)
            task.status = ExportStatus.COMPLETED
            task.progress = 100
            task.completed_at = datetime.utcnow()
            task.expires_at = datetime.utcnow() + timedelta(hours=24)

        except Exception as e:
            task.status = ExportStatus.FAILED
            task.error_message = str(e)
            task.progress = 0

    async def _get_presentation(
        self, presentation_id: UUID
    ) -> Optional["Presentation"]:
        """获取演示文稿"""
        from ai_ppt.domain.models.presentation import Presentation

        stmt = select(Presentation).where(Presentation.id == presentation_id)
        result = await self._session.execute(stmt)
        return result.scalar_one_or_none()

    async def _get_slides(self, presentation_id: UUID) -> list["Slide"]:
        """获取幻灯片列表"""
        from ai_ppt.domain.models.slide import Slide

        stmt = (
            select(Slide)
            .where(Slide.presentation_id == presentation_id)
            .order_by(Slide.order_index)
        )
        result = await self._session.execute(stmt)
        return list(result.scalars().all())

    def _filter_slides_by_range(self, slides: list, slide_range: str) -> list:
        """根据范围过滤幻灯片"""
        try:
            if "-" in slide_range:
                start, end = map(int, slide_range.split("-"))
                return slides[start - 1 : end]
            else:
                idx = int(slide_range) - 1
                return [slides[idx]] if 0 <= idx < len(slides) else slides
        except (ValueError, IndexError):
            return slides

    async def _export_pptx(
        self,
        presentation: "Presentation",
        slides: list["Slide"],
        task: ExportTask,
    ) -> str:
        """
        导出为 PPTX 格式

        Args:
            presentation: 演示文稿
            slides: 幻灯片列表
            task: 导出任务

        Returns:
            文件路径
        """
        from pptx import Presentation as PptxPresentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt

        task.progress = 60

        # 创建 PPTX 演示文稿
        prs = PptxPresentation()

        # 设置幻灯片尺寸 (16:9)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        task.progress = 70

        # 根据主题设置背景色
        theme_colors = self._get_theme_colors(presentation.theme)

        for idx, slide in enumerate(slides):
            # 添加幻灯片
            blank_layout = prs.slide_layouts[6]  # 空白布局
            pptx_slide = prs.slides.add_slide(blank_layout)

            # 设置背景
            background = pptx_slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*theme_colors["background"])

            # 添加内容
            content = slide.content or {}

            # 标题
            if content.get("title"):
                title_box = pptx_slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
                )
                tf = title_box.text_frame
                p = tf.paragraphs[0]
                p.text = content["title"]
                p.font.size = Pt(32 if idx == 0 else 28)
                p.font.bold = True
                p.font.color.rgb = RGBColor(*theme_colors["title"])

            # 副标题
            if content.get("subtitle"):
                subtitle_box = pptx_slide.shapes.add_textbox(
                    Inches(0.5), Inches(1.2), Inches(12.333), Inches(0.5)
                )
                tf = subtitle_box.text_frame
                p = tf.paragraphs[0]
                p.text = content["subtitle"]
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(*theme_colors["text"])

            # 正文内容
            if content.get("text"):
                body_box = pptx_slide.shapes.add_textbox(
                    Inches(0.5), Inches(2), Inches(12.333), Inches(4.5)
                )
                tf = body_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = content["text"]
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(*theme_colors["text"])

            # 项目符号列表
            if content.get("bullets"):
                bullets_box = pptx_slide.shapes.add_textbox(
                    Inches(0.5), Inches(2), Inches(12.333), Inches(4.5)
                )
                tf = bullets_box.text_frame
                tf.word_wrap = True
                for i, bullet in enumerate(content["bullets"]):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = f"• {bullet}"
                    p.font.size = Pt(14)
                    p.font.color.rgb = RGBColor(*theme_colors["text"])
                    p.space_after = Pt(8)

            # 备注
            if task.include_notes and slide.notes:
                pptx_slide.notes_slide.notes_text_frame.text = slide.notes

            task.progress = 70 + int((idx + 1) / len(slides) * 20)

        task.progress = 95

        # 保存文件
        filename = f"{presentation.id}_{uuid.uuid4().hex[:8]}.pptx"
        file_path = self._exports_dir / filename
        prs.save(str(file_path))

        return str(file_path)

    async def _export_pdf(
        self,
        presentation: "Presentation",
        slides: list["Slide"],
        task: ExportTask,
    ) -> str:
        """
        导出为 PDF 格式

        Args:
            presentation: 演示文稿
            slides: 幻灯片列表
            task: 导出任务

        Returns:
            文件路径
        """
        from reportlab.lib.pagesizes import A4, landscape
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        from reportlab.pdfgen import canvas

        task.progress = 60

        # 注册中文字体（如果可用）
        try:
            pdfmetrics.registerFont(
                TTFont("SimHei", "/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc")
            )
            font_name = "SimHei"
        except Exception:
            font_name = "Helvetica"

        # 创建 PDF
        filename = f"{presentation.id}_{uuid.uuid4().hex[:8]}.pdf"
        file_path = self._exports_dir / filename

        # 设置页面尺寸 (16:9)
        width, height = landscape(A4)

        c = canvas.Canvas(str(file_path), pagesize=(width, height))

        task.progress = 70

        theme_colors = self._get_theme_colors(presentation.theme)

        for idx, slide in enumerate(slides):
            # 背景色
            bg_color = theme_colors["background"]
            c.setFillColorRGB(bg_color[0] / 255, bg_color[1] / 255, bg_color[2] / 255)
            c.rect(0, 0, width, height, fill=1, stroke=0)

            content = slide.content or {}

            # 标题
            if content.get("title"):
                title_color = theme_colors["title"]
                c.setFillColorRGB(
                    title_color[0] / 255, title_color[1] / 255, title_color[2] / 255
                )
                c.setFont(font_name, 28 if idx == 0 else 24)
                c.drawString(40, height - 60, content["title"])

            # 副标题
            if content.get("subtitle"):
                text_color = theme_colors["text"]
                c.setFillColorRGB(
                    text_color[0] / 255, text_color[1] / 255, text_color[2] / 255
                )
                c.setFont(font_name, 16)
                c.drawString(40, height - 100, content["subtitle"])

            # 正文
            y_pos = height - 150
            if content.get("text"):
                text_color = theme_colors["text"]
                c.setFillColorRGB(
                    text_color[0] / 255, text_color[1] / 255, text_color[2] / 255
                )
                c.setFont(font_name, 12)

                # 简单的文本换行
                text = content["text"]
                words = text.split()
                line = ""
                for word in words:
                    if c.stringWidth(line + " " + word, font_name, 12) < width - 80:
                        line += " " + word if line else word
                    else:
                        c.drawString(40, y_pos, line)
                        y_pos -= 20
                        line = word
                if line:
                    c.drawString(40, y_pos, line)

            # 项目符号
            if content.get("bullets"):
                text_color = theme_colors["text"]
                c.setFillColorRGB(
                    text_color[0] / 255, text_color[1] / 255, text_color[2] / 255
                )
                c.setFont(font_name, 12)

                for bullet in content["bullets"]:
                    c.drawString(40, y_pos, f"• {bullet}")
                    y_pos -= 25

            # 备注
            if task.include_notes and slide.notes:
                c.setFont(font_name, 10)
                c.setFillColorRGB(0.5, 0.5, 0.5)
                c.drawString(40, 40, f"Notes: {slide.notes[:100]}...")

            c.showPage()
            task.progress = 70 + int((idx + 1) / len(slides) * 20)

        task.progress = 95
        c.save()

        return str(file_path)

    async def _export_images(
        self,
        presentation: "Presentation",
        slides: list["Slide"],
        task: ExportTask,
    ) -> str:
        """
        导出为图片格式

        Args:
            presentation: 演示文稿
            slides: 幻灯片列表
            task: 导出任务

        Returns:
            文件路径（压缩包）
        """
        import zipfile

        task.progress = 60

        # 设置分辨率
        if task.quality == "high":
            width, height = 1920, 1080
        else:
            width, height = 1280, 720

        theme_colors = self._get_theme_colors(presentation.theme)

        # 创建临时目录存储图片
        temp_dir = self._exports_dir / f"temp_{uuid.uuid4().hex[:8]}"
        temp_dir.mkdir(exist_ok=True)

        image_files = []

        for idx, slide in enumerate(slides):
            # 创建图片
            img = Image.new("RGB", (width, height), theme_colors["background"])
            draw = ImageDraw.Draw(img)

            # 尝试加载字体
            title_font: Union[ImageFont.FreeTypeFont, ImageFont.ImageFont]
            subtitle_font: Union[ImageFont.FreeTypeFont, ImageFont.ImageFont]
            text_font: Union[ImageFont.FreeTypeFont, ImageFont.ImageFont]
            try:
                title_font = ImageFont.truetype(
                    "/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc",
                    48 if idx == 0 else 40,
                )
                subtitle_font = ImageFont.truetype(
                    "/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc", 28
                )
                text_font = ImageFont.truetype(
                    "/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc", 20
                )
            except Exception:
                title_font = ImageFont.load_default()
                subtitle_font = ImageFont.load_default()
                text_font = ImageFont.load_default()

            content = slide.content or {}

            # 标题
            if content.get("title"):
                draw.text(
                    (40, 30),
                    content["title"],
                    fill=theme_colors["title"],
                    font=title_font,
                )

            # 副标题
            if content.get("subtitle"):
                draw.text(
                    (40, 100),
                    content["subtitle"],
                    fill=theme_colors["text"],
                    font=subtitle_font,
                )

            # 正文
            y_pos = 180
            if content.get("text"):
                # 简单的文本换行
                text = content["text"]
                words = text.split()
                line = ""
                for word in words:
                    bbox = draw.textbbox((0, 0), line + " " + word, font=text_font)
                    if bbox[2] < width - 80:
                        line += " " + word if line else word
                    else:
                        draw.text(
                            (40, y_pos), line, fill=theme_colors["text"], font=text_font
                        )
                        y_pos += 30
                        line = word
                if line:
                    draw.text(
                        (40, y_pos), line, fill=theme_colors["text"], font=text_font
                    )

            # 项目符号
            if content.get("bullets"):
                y_pos = max(y_pos + 40, 180)
                for bullet in content["bullets"]:
                    draw.text(
                        (40, y_pos),
                        f"• {bullet}",
                        fill=theme_colors["text"],
                        font=text_font,
                    )
                    y_pos += 35

            # 保存图片
            ext = "png" if task.format == ExportFormat.PNG else "jpg"
            img_filename = f"slide_{idx + 1:03d}.{ext}"
            img_path = temp_dir / img_filename

            if task.format == ExportFormat.PNG:
                img.save(img_path, "PNG")
            else:
                img.save(img_path, "JPEG", quality=95 if task.quality == "high" else 80)

            image_files.append(img_path)
            task.progress = 70 + int((idx + 1) / len(slides) * 15)

        task.progress = 90

        # 打包为 zip
        zip_filename = f"{presentation.id}_{uuid.uuid4().hex[:8]}.zip"
        zip_path = self._exports_dir / zip_filename

        with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
            for img_file in image_files:
                zf.write(img_file, img_file.name)

        # 清理临时文件
        for img_file in image_files:
            img_file.unlink()
        temp_dir.rmdir()

        task.progress = 95

        return str(zip_path)

    def _get_theme_colors(self, theme: str) -> dict:
        """
        获取主题颜色

        Args:
            theme: 主题名称

        Returns:
            颜色配置字典
        """
        themes = {
            "default": {
                "background": (255, 255, 255),
                "title": (51, 51, 51),
                "text": (102, 102, 102),
            },
            "dark": {
                "background": (45, 45, 45),
                "title": (255, 255, 255),
                "text": (200, 200, 200),
            },
            "blue": {
                "background": (240, 248, 255),
                "title": (25, 55, 109),
                "text": (60, 80, 120),
            },
            "green": {
                "background": (240, 255, 240),
                "title": (34, 85, 51),
                "text": (60, 100, 70),
            },
        }
        return themes.get(theme, themes["default"])

    def get_full_path(self, file_path: str) -> Path:
        """
        获取完整文件路径

        Args:
            file_path: 相对或绝对路径

        Returns:
            完整路径
        """
        path = Path(file_path)
        if path.is_absolute():
            return path
        return self._exports_dir / path

    def get_file_url(self, file_path: str) -> str:
        """
        获取文件下载 URL

        Args:
            file_path: 文件路径

        Returns:
            下载 URL
        """
        filename = Path(file_path).name
        return f"/api/v1/exports/download/{filename}"


# 后台任务处理
async def process_export_task(task_id: UUID) -> None:
    """
    处理导出任务（后台执行）

    Args:
        task_id: 任务ID
    """
    # 创建新的会话处理任务
    from ai_ppt.database import AsyncSessionLocal

    async with AsyncSessionLocal() as session:
        service = ExportService(session)
        await service.process_export(task_id)
